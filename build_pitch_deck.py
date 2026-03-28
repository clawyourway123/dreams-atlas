#!/usr/bin/env python3
"""Build K-Dense Investor Pitch Deck PPTX with professional design."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === Design Constants ===
PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)      # Deep navy
ACCENT = RGBColor(0x2D, 0x8B, 0x7E)       # Teal
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
TEXT_BODY = RGBColor(0x44, 0x44, 0x44)
TABLE_HEADER_BG = PRIMARY
TABLE_ALT_ROW = RGBColor(0xF0, 0xF5, 0xF8)
TABLE_BORDER = RGBColor(0xCC, 0xCC, 0xCC)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_NAME = "Calibri"

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_background(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_navy_bar(slide, top=0, height=Inches(1.2)):
    """Add a navy header bar across the top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    return shape


def add_accent_line(slide, top, left=Inches(0.8), width=Inches(2)):
    """Add a teal accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def add_footer(slide, text="K-Dense Science Lab  |  DREAMS Project  |  Confidential"):
    """Add a consistent footer."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_HEIGHT - Inches(0.45), SLIDE_WIDTH, Inches(0.45)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), SLIDE_HEIGHT - Inches(0.4), Inches(11), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.font.name = FONT_NAME


def add_slide_title(slide, title, subtitle=None):
    """Add title bar with optional subtitle."""
    add_navy_bar(slide, height=Inches(1.1))
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11), Inches(0.35))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.italic = True
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
        p2.font.name = FONT_NAME


def add_body_text(slide, text, left=Inches(0.8), top=Inches(1.5), width=Inches(11.5), height=Inches(5),
                  font_size=Pt(16), bold=False, color=TEXT_BODY):
    """Add body text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.space_after = Pt(6)
    return tf


def add_bullet_list(slide, items, left=Inches(0.8), top=Inches(1.5), width=Inches(11.5), height=Inches(5),
                    font_size=Pt(15), color=TEXT_BODY):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(8)
        p.level = 0
    return tf


def add_table(slide, data, left=Inches(0.8), top=Inches(2.5), width=Inches(11.5), row_height=Inches(0.4)):
    """Add a professionally formatted table."""
    rows = len(data)
    cols = len(data[0])
    col_width = int(width / cols)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, row_height * rows)
    table = table_shape.table

    for col_idx in range(cols):
        table.columns[col_idx].width = col_width

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = FONT_NAME

                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_GRAY
                    if col_idx == 0:
                        paragraph.font.bold = True

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Cell fill
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_ROW
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


def add_highlight_box(slide, text, left, top, width, height, bg_color=ACCENT, text_color=WHITE, font_size=Pt(14)):
    """Add a colored highlight box with text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = text_color
    p.font.name = FONT_NAME
    p.font.bold = True
    return shape


# ============================================================
# SLIDE 1: Title / Company Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, PRIMARY)

# Large company name
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "K-Dense Science Lab"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = FONT_NAME
p.alignment = PP_ALIGN.CENTER

# Tagline
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(0.6))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "AI-Powered Material Identification Through Spectral Intelligence"
p2.font.size = Pt(22)
p2.font.italic = True
p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
p2.font.name = FONT_NAME
p2.alignment = PP_ALIGN.CENTER

# Accent line
add_accent_line(slide, top=Inches(3.6), left=Inches(5), width=Inches(3))

# Key points
bullets = [
    "Deep research laboratory combining AI/ML with analytical chemistry",
    "Cross-disciplinary team: machine learning, spectroscopy, data science",
    "Flagship: Automated adhesive classification using IR and Raman spectroscopy",
]
txBox3 = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9), Inches(2.5))
tf3 = txBox3.text_frame
tf3.word_wrap = True
for i, b in enumerate(bullets):
    if i == 0:
        p = tf3.paragraphs[0]
    else:
        p = tf3.add_paragraph()
    p.text = f"\u2022  {b}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

# Footer - DREAMS project
txBox4 = slide.shapes.add_textbox(Inches(1), Inches(6.6), Inches(11), Inches(0.5))
tf4 = txBox4.text_frame
p4 = tf4.paragraphs[0]
p4.text = "DREAMS Project  |  Investor Presentation  |  2026"
p4.font.size = Pt(12)
p4.font.color.rgb = RGBColor(0x88, 0x99, 0xAA)
p4.font.name = FONT_NAME
p4.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_slide_title(slide, "The Problem", "Adhesive Testing Is Broken")
add_footer(slide)

bullets = [
    "Traditional adhesive identification relies on slow, expensive wet-lab methods (ASTM/ISO protocols)",
    "Testing a single adhesive sample can take hours to days and costs $50\u2013200+ per test",
    "Current methods require specialized chemists and hazardous reagents",
    "Quality control bottlenecks cause production delays in aerospace, automotive, electronics, and packaging",
    "Adhesive misidentification leads to product failures, recalls, and liability",
    "No scalable, automated solution exists for real-time adhesive classification",
]
add_bullet_list(slide, [f"\u2022  {b}" for b in bullets], top=Inches(1.5), font_size=Pt(16))

add_highlight_box(slide, "Global adhesive market exceeds $65B \u2014 No automated classification solution exists",
                  Inches(1.5), Inches(5.8), Inches(10), Inches(0.6), bg_color=PRIMARY)

# ============================================================
# SLIDE 3: Our Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_slide_title(slide, "Our Solution", "AI-Powered Spectral Classification \u2014 Instant, Accurate, Non-Destructive")
add_footer(slide)

solution_points = [
    "Combine IR (infrared) and Raman spectroscopy with trained ML models",
    "Classify adhesives into 7 major categories in seconds, not hours",
    "Non-destructive testing \u2014 no sample preparation, no reagents",
    "Works with existing benchtop and portable IR/Raman spectrometers",
    "Cloud or on-premise deployment for real-time QC integration",
]
add_bullet_list(slide, [f"\u2022  {b}" for b in solution_points], top=Inches(1.5), font_size=Pt(16))

# Categories boxes
categories = ["Acrylic/PSA", "Cyanoacrylate", "Epoxy", "Hot-melt", "Polyurethane", "Rubber-based", "Silicone"]
x_start = Inches(0.8)
for i, cat in enumerate(categories):
    add_highlight_box(slide, cat, x_start + Inches(i * 1.75), Inches(5.5), Inches(1.6), Inches(0.55),
                      bg_color=ACCENT, font_size=Pt(11))

# ============================================================
# SLIDE 4: Technology \u2014 Validated Performance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_slide_title(slide, "Technology \u2014 Validated Performance", "IR + Raman Achieves Near-Perfect Classification")
add_footer(slide)

perf_data = [
    ["Metric", "IR/FTIR Alone", "Raman Alone", "Combined IR+Raman"],
    ["RF Accuracy", "100.0%", "100.0%", "95\u2013100%"],
    ["CNN-1D Accuracy", "95.2%", "99.8%", "95%+"],
    ["PLS-DA Accuracy", "100.0%", "100.0%", "100%"],
]
add_table(slide, perf_data, top=Inches(1.5), width=Inches(11.5))

# Key differentiators
diff_points = [
    "Compound-grouped 5-fold cross-validation (no data leakage)",
    "Domain-validated: model features align with known adhesive chemistry",
    "Fisher ratios > 23 for both IR and Raman (excellent class separation)",
    "Dataset: 1,500+ curated adhesive spectra across 7 classes",
]
add_bullet_list(slide, [f"\u2022  {d}" for d in diff_points], top=Inches(3.8), font_size=Pt(14))

# Chemistry callouts
chem_points = [
    "Epoxide ring breathing at 910\u2013915 cm\u207b\u00b9",
    "Urethane N-H stretch at 3300 cm\u207b\u00b9",
    "Silicone Si-O-Si at 1000\u20131100 cm\u207b\u00b9",
]
txBox = slide.shapes.add_textbox(Inches(7), Inches(5.2), Inches(5), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, cp in enumerate(chem_points):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = cp
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = ACCENT
    p.font.name = FONT_NAME

# ============================================================
# SLIDE 5: Market Opportunity
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_slide_title(slide, "Market Opportunity", "$420M+ TAM \u2014 No Direct Competitor")
add_footer(slide)

market_data = [
    ["Market Layer", "Value"],
    ["Total Addressable Market (TAM)", "~$500\u2013620M"],
    ["Serviceable Addressable Market (SAM)", "~$165M"],
    ["Serviceable Obtainable Market (SOM, Year 5)", "~$12M ARR"],
]
add_table(slide, market_data, top=Inches(1.5), width=Inches(8))

# Bottom-up construction
bu_data = [
    ["Tier", "Companies", "Avg Deal", "Total"],
    ["Tier 1 (>$10B rev)", "~20", "$800K", "$16M"],
    ["Tier 2 ($1B\u2013$10B)", "~180", "$400K", "$72M"],
    ["Tier 3 ($500M\u2013$1B)", "~1,900", "$175K", "$332M"],
]
add_table(slide, bu_data, top=Inches(3.5), width=Inches(8))

# Market tailwinds
tailwinds = [
    "PFAS reformulation wave driving urgent need",
    "EU Digital Product Passport regulations",
    "AI-assisted formulation trend (5.7% CAGR)",
]
txBox = slide.shapes.add_textbox(Inches(8.5), Inches(1.5), Inches(4.5), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Market Tailwinds"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.font.name = FONT_NAME
for tw in tailwinds:
    p = tf.add_paragraph()
    p.text = f"\u2713  {tw}"
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT
    p.font.name = FONT_NAME
    p.space_after = Pt(8)

add_highlight_box(slide, "Category Creator \u2014 $1.2M+ barrier to replicate dataset & models",
                  Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.55), bg_color=PRIMARY)

# ============================================================
# SLIDE 6: Business Model & Pricing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_slide_title(slide, "Business Model & Pricing", "SaaS-First, Multi-Tier Revenue Model")
add_footer(slide)

pricing_data = [
    ["Tier", "Price", "Target", "% Revenue"],
    ["SaaS Platform", "$2,500\u20138,000/mo", "R&D teams", "60%"],
    ["Per-Test API", "$75\u2013150/test", "QC labs", "25%"],
    ["Enterprise License", "$120\u2013250K/yr", "On-prem", "15%"],
    ["Pilot Program", "$50K flat (6 mo)", "New customers", "Entry"],
    ["Academic", "$5\u201315K/yr", "Universities", "Self-serve"],
]
add_table(slide, pricing_data, top=Inches(1.4), width=Inches(11.5))

rev_data = [
    ["Year", "Customers", "ARR", "YoY Growth"],
    ["2026", "24", "$727K", "\u2014"],
    ["2027", "60", "$2.2M", "199%"],
    ["2028", "100", "$4.0M", "85%"],
    ["2029", "143", "$6.4M", "59%"],
    ["2030", "187", "$8.7M", "37%"],
]
add_table(slide, rev_data, top=Inches(4.2), width=Inches(8))

add_highlight_box(slide, "Per-test upside: 100K tests/yr across 50 QC labs = $10\u201325M opportunity",
                  Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5), bg_color=ACCENT, font_size=Pt(13))

# ============================================================
# SLIDE 7: Financial Projections
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_slide_title(slide, "Financial Projections", "Path to $8.7M ARR and 45% Net Margins by Year 5")
add_footer(slide)

fin_data = [
    ["Metric", "Year 1", "Year 2", "Year 3", "Year 4", "Year 5"],
    ["Revenue", "$727K", "$2.2M", "$4.0M", "$6.4M", "$8.7M"],
    ["Gross Margin", "58%", "68%", "74%", "77%", "78%"],
    ["EBITDA", "($523K)", "$274K", "$1.5M", "$3.2M", "$5.0M"],
    ["Net Margin", "(77%)", "8%", "27%", "39%", "45%"],
    ["Customers", "24", "60", "100", "143", "187"],
]
add_table(slide, fin_data, top=Inches(1.4), width=Inches(11.5))

unit_data = [
    ["Metric", "Year 1", "Year 5"],
    ["CAC (fully loaded)", "$18.8K", "$8.5K"],
    ["LTV (3-year, blended)", "$62K", "$102K"],
    ["LTV:CAC", "3.3x", "12.0x"],
    ["Payback period", "14 months", "6 months"],
    ["Annual churn", "15%", "7%"],
]
add_table(slide, unit_data, top=Inches(4.2), width=Inches(7))

add_highlight_box(slide, "Break-even: Month 26 (~55 customers, $210K MRR)",
                  Inches(7.5), Inches(4.5), Inches(4.8), Inches(0.6), bg_color=PRIMARY)

# ============================================================
# SLIDE 8: Competitive Advantages
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_slide_title(slide, "Competitive Advantages", "Why K-Dense Wins")
add_footer(slide)

comp_data = [
    ["Advantage", "K-Dense", "Traditional Labs", "DIY (Python)"],
    ["Speed", "Seconds", "Hours\u2013Days", "Minutes"],
    ["Accuracy", "95\u2013100%", "Variable", "Untested"],
    ["Sample Prep", "None", "Extensive", "N/A"],
    ["Cost per Test", "Low (SaaS)", "$50\u2013200+", "Engineer time"],
    ["Non-Destructive", "Yes", "No", "N/A"],
    ["7-Class Coverage", "Yes", "Yes", "Limited"],
]
add_table(slide, comp_data, top=Inches(1.4), width=Inches(11.5))

moats = [
    "Proprietary dataset \u2014 1,500+ labeled adhesive spectra ($1.2M+ to replicate)",
    "Domain-validated models \u2014 chemically interpretable features",
    "Compound-grouped validation \u2014 exceeds industry standards",
    "Customer data moat \u2014 switching costs grow over time",
    "First-mover advantage \u2014 no direct competitor",
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Defensible Moats"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.font.name = FONT_NAME
p.space_after = Pt(6)
for i, moat in enumerate(moats):
    p = tf.add_paragraph()
    p.text = f"{i+1}.  {moat}"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_BODY
    p.font.name = FONT_NAME
    p.space_after = Pt(4)

# ============================================================
# SLIDE 9: Team & Capabilities
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_slide_title(slide, "Team & Capabilities", "World-Class Interdisciplinary Team")
add_footer(slide)

team_members = [
    ("Dr. Elena Vasquez", "CEO", "Strategic vision and company leadership"),
    ("Dr. Marcus Chen", "Chief Science Officer", "Research coordination"),
    ("Dr. Alexander Petrov", "Machine Learning", "Model architecture, training, production deployment"),
    ("Dr. Nikolai Volkov", "Physical Sciences", "Spectroscopy and domain science"),
    ("Dr. Priya Sharma", "Chemistry", "Domain science and validation"),
    ("Dr. Lin Wei", "Data Engineering", "Data curation, pipeline, quality assurance"),
    ("Dr. Rosa Martinez", "Visualization", "Data visualization and communication"),
    ("Dr. Victoria Chang", "Communication", "Scientific communication lead"),
    ("Dr. Paolo Ricci", "Market Analysis", "Market research and business strategy"),
    ("Dr. Sarah Nakamura", "Finance", "Financial modeling and projections"),
]

# Two columns
col1_x = Inches(0.8)
col2_x = Inches(6.8)
y_start = Inches(1.5)
row_h = Inches(0.5)

for i, (name, role, desc) in enumerate(team_members):
    col_x = col1_x if i < 5 else col2_x
    row_y = y_start + (i % 5) * row_h

    txBox = slide.shapes.add_textbox(col_x, row_y, Inches(5.5), Inches(0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{name}  "
    run1.font.size = Pt(14)
    run1.font.bold = True
    run1.font.color.rgb = PRIMARY
    run1.font.name = FONT_NAME
    run2 = p.add_run()
    run2.text = f"({role}) \u2014 {desc}"
    run2.font.size = Pt(12)
    run2.font.color.rgb = TEXT_BODY
    run2.font.name = FONT_NAME

add_highlight_box(slide, "Full team of 15+ specialists across ML, chemistry, data science, and business strategy",
                  Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.55), bg_color=ACCENT, font_size=Pt(14))

# ============================================================
# SLIDE 10: Go-to-Market Strategy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_slide_title(slide, "Go-to-Market Strategy", "Land-and-Expand with Design Partners")
add_footer(slide)

partner_data = [
    ["Target", "Pilot Probability", "Rationale"],
    ["Evonik Industries", "88%", "Open Innovation unit, digital chemistry commitment"],
    ["Henkel AG", "82%", "World's largest adhesive manufacturer, active Digital R&D"],
    ["Covestro AG", "75%", "Strong digital culture, PU adhesive formulation focus"],
    ["Arkema/Bostik", "70%", "Post-acquisition data integration need"],
    ["3M Company", "64%", "Largest adhesive portfolio (long procurement cycle)"],
]
add_table(slide, partner_data, top=Inches(1.4), width=Inches(11.5))

phase_data = [
    ["Phase", "Timeline", "Goal", "Revenue Target"],
    ["Foundation", "Q1\u2013Q2 2026", "SOC 2, vendor portals, V1 deployment", "\u2014"],
    ["Design Partners", "Q3\u2013Q4 2026", "Sign 2\u20133 paid pilots at $50K each", "$100\u2013150K"],
    ["Enterprise Expansion", "2027", "Convert pilots to $200\u2013300K contracts", "$1.5\u20132.5M ARR"],
    ["Scale", "2028\u20132030", "10\u201315 enterprise + mid-tier", "$5\u201312M ARR"],
]
add_table(slide, phase_data, top=Inches(4.2), width=Inches(11.5))

# ============================================================
# SLIDE 11: Investment Ask
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_slide_title(slide, "Investment Ask", "$1.5M Seed \u2192 $5M Series A")
add_footer(slide)

# Seed round
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Seed Round: $1.5M"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.font.name = FONT_NAME

seed_data = [
    ["Use of Funds", "Amount"],
    ["Product development & model optimization", "$600K"],
    ["Initial sales team (2 AEs + SDR)", "$350K"],
    ["Cloud infrastructure & hosting", "$200K"],
    ["Regulatory & compliance (SOC 2, ISO)", "$100K"],
    ["Working capital", "$250K"],
]
add_table(slide, seed_data, top=Inches(1.9), width=Inches(5.5))

# Series A
txBox2 = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.5), Inches(0.4))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Series A: $5M (Month 18\u201324)"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.font.name = FONT_NAME

series_a_data = [
    ["Use of Funds", "Amount"],
    ["Scale engineering (8\u219212 FTEs)", "$2.0M"],
    ["Expand sales & marketing", "$1.5M"],
    ["Enterprise on-prem features", "$800K"],
    ["International market entry", "$500K"],
    ["Working capital", "$200K"],
]
add_table(slide, series_a_data, left=Inches(7), top=Inches(1.9), width=Inches(5.5))

# Milestones
milestones_seed = [
    "Launch commercial SaaS product (Month 6)",
    "15+ paying customers (Month 12)",
    "$50K MRR (Month 15)",
    "Validate enterprise licensing (Month 18)",
]
milestones_a = [
    "$200K MRR (Month 24)",
    "60+ customers (Month 30)",
    "Positive EBITDA (Month 26)",
    "First on-prem deployment (Month 24)",
]

txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(5.5), Inches(2))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "Seed Milestones (18 months)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = FONT_NAME
for m in milestones_seed:
    p = tf3.add_paragraph()
    p.text = f"\u2713  {m}"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_BODY
    p.font.name = FONT_NAME

txBox4 = slide.shapes.add_textbox(Inches(7), Inches(4.8), Inches(5.5), Inches(2))
tf4 = txBox4.text_frame
tf4.word_wrap = True
p = tf4.paragraphs[0]
p.text = "Series A Milestones"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = FONT_NAME
for m in milestones_a:
    p = tf4.add_paragraph()
    p.text = f"\u2713  {m}"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_BODY
    p.font.name = FONT_NAME

# ============================================================
# SLIDE 12: Vision / Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, PRIMARY)

# Quote
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"Every adhesive bond in every product deserves to be verified \u2014 instantly, accurately, and affordably."'
p.font.size = Pt(22)
p.font.italic = True
p.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
p.font.name = FONT_NAME
p.alignment = PP_ALIGN.CENTER

add_accent_line(slide, top=Inches(2.5), left=Inches(5), width=Inches(3))

# Vision points
vision = [
    ("Near-term", "Production-ready adhesive classifier with 95\u2013100% accuracy"),
    ("Medium-term", "Industry-standard AI testing platform integrated with major spectrometer brands"),
    ("Long-term", "Universal spectral intelligence platform for material identification across industries"),
]
y = Inches(3.0)
for label, desc in vision:
    txBox = slide.shapes.add_textbox(Inches(1.5), y, Inches(10), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{label}:  "
    run1.font.size = Pt(16)
    run1.font.bold = True
    run1.font.color.rgb = ACCENT
    run1.font.name = FONT_NAME
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(16)
    run2.font.color.rgb = WHITE
    run2.font.name = FONT_NAME
    y += Inches(0.55)

# Scenario table
scenario_data = [
    ["Scenario", "Year 5 ARR", "Assumption"],
    ["Bear", "$1.5M", "Pilot struggles; build-in-house risk"],
    ["Base", "$8.7M", "3 pilots in 2026; steady expansion"],
    ["Bull", "$14M", "PFAS mandate urgency + adjacent verticals"],
]
add_table(slide, scenario_data, top=Inches(5.0), width=Inches(10), left=Inches(1.5))

# Contact
txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "K-Dense Science Lab  \u2014  DREAMS Project  \u2014  Contact us to schedule a demo"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
p.font.name = FONT_NAME
p.alignment = PP_ALIGN.CENTER


# === Save ===
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "K-Dense_Investor_Pitch_Deck.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
