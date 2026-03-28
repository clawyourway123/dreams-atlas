#!/usr/bin/env python3
"""Generate investor pitch deck PPTX for K-Dense Science Lab DREAMS Project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)       # Deep navy
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)    # K-Dense blue
ACCENT_GREEN = RGBColor(0x00, 0xB8, 0x8D)   # Teal/green accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GOLD = RGBColor(0xFF, 0xC1, 0x07)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
SUBTITLE_GRAY = RGBColor(0x88, 0x88, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BG):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def tf_para(tf, text, size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, space_after=Pt(6), space_before=Pt(0)):
    """Add a paragraph to a text frame."""
    p = tf.add_paragraph() if len(tf.paragraphs) > 0 and tf.paragraphs[0].text != '' else tf.paragraphs[0]
    if len(tf.paragraphs) > 1 or tf.paragraphs[0].text != '':
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    return p


def add_text_box(slide, left, top, width, height):
    """Add a text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_bullet(tf, text, size=16, color=WHITE, level=0, bold=False):
    """Add a bullet point paragraph."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_after = Pt(4)
    return p


def make_table(slide, rows_data, left, top, width, row_height=Inches(0.4), header_color=ACCENT_BLUE, font_size=12):
    """Create a styled table."""
    rows = len(rows_data)
    cols = len(rows_data[0])
    col_width = width // cols

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, row_height * rows)
    table = table_shape.table

    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = WHITE if i == 0 else RGBColor(0xDD, 0xDD, 0xDD)
                paragraph.font.bold = (i == 0)
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            if i == 0:
                cell.fill.fore_color.rgb = header_color
            elif i % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x3A)
            else:
                cell.fill.fore_color.rgb = RGBColor(0x15, 0x1E, 0x32)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Accent bar at top
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

# Company name
tf = add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2))
tf_para(tf, "K-DENSE SCIENCE LAB", size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Subtitle
tf = add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8))
tf_para(tf, "AI-Powered Material Identification Through Spectral Intelligence", size=24, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

# Divider
add_shape_bg(slide, Inches(5), Inches(3.9), Inches(3), Inches(0.04), ACCENT_BLUE)

# Deck title
tf = add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(1))
tf_para(tf, "Investor Pitch Deck", size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
tf_para(tf, "DREAMS Project  |  March 2026", size=18, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

# Key stat
tf = add_text_box(slide, Inches(2), Inches(5.8), Inches(9), Inches(0.6))
tf_para(tf, "95-100% Accuracy  \u2022  7 Adhesive Classes  \u2022  Non-Destructive  \u2022  Real-Time", size=16, color=GOLD, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
tf_para(tf, "THE PROBLEM", size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.LEFT)
tf_para(tf, "Adhesive Testing Is Broken", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

tf = add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
bullets = [
    ("Traditional adhesive identification relies on slow, expensive wet-lab methods (ASTM/ISO protocols)", False),
    ("Testing a single sample: hours to days, $50\u2013200+ per test", False),
    ("Requires specialized chemists and hazardous reagents", False),
    ("QC bottlenecks cause production delays in aerospace, automotive, electronics, and packaging", False),
    ("Adhesive misidentification leads to product failures, recalls, and liability", False),
    ("Global adhesive market exceeds $65B \u2014 no scalable, automated classification solution exists", True),
]
for text, bold in bullets:
    add_bullet(tf, "\u2022  " + text, size=18, color=WHITE if not bold else GOLD, bold=bold)

# Key stat box
box = add_shape_bg(slide, Inches(3), Inches(5.8), Inches(7), Inches(1.2), RGBColor(0x1A, 0x25, 0x3A))
tf2 = box.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "$65B+ market  \u2022  No automated solution  \u2022  Hours per test  \u2022  $200+ per sample"
p.font.size = Pt(16)
p.font.color.rgb = GOLD
p.font.bold = True
p.alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 3: OUR SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_GREEN)

tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
tf_para(tf, "OUR SOLUTION", size=14, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.LEFT)
tf_para(tf, "AI-Powered Spectral Classification", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

tf = add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
add_bullet(tf, "\u2022  Combine IR and Raman spectroscopy with trained ML models", size=17, color=WHITE)
add_bullet(tf, "\u2022  Classify adhesives into 7 major categories in seconds", size=17, color=WHITE)
add_bullet(tf, "\u2022  Non-destructive \u2014 no sample prep, no reagents", size=17, color=WHITE)
add_bullet(tf, "\u2022  Works with existing benchtop and portable spectrometers", size=17, color=WHITE)
add_bullet(tf, "\u2022  Cloud or on-premise deployment", size=17, color=WHITE)

# 7 classes box
box = add_shape_bg(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5), RGBColor(0x1A, 0x25, 0x3A))
tf2 = box.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "7 Adhesive Classes"
p.font.size = Pt(20)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

classes = ["Acrylic / PSA", "Cyanoacrylate", "Epoxy", "Hot-melt", "Polyurethane", "Rubber-based", "Silicone"]
for cls in classes:
    pp = tf2.add_paragraph()
    pp.text = "\u2713  " + cls
    pp.font.size = Pt(16)
    pp.font.color.rgb = WHITE
    pp.alignment = PP_ALIGN.LEFT
    pp.space_before = Pt(6)

# Bottom highlight
tf = add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8))
tf_para(tf, "Instant  \u2022  Accurate  \u2022  Non-Destructive  \u2022  Affordable", size=22, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: TECHNOLOGY \u2014 VALIDATED PERFORMANCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
tf_para(tf, "TECHNOLOGY", size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.LEFT)
tf_para(tf, "Validated Performance: 95\u2013100% Accuracy", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

# Performance table
table_data = [
    ["Model", "IR/FTIR Alone", "Raman Alone", "Combined IR+Raman"],
    ["Random Forest", "100.0%", "100.0%", "95\u2013100%"],
    ["CNN-1D", "95.2%", "99.8%", "95%+"],
    ["PLS-DA", "100.0%", "100.0%", "100%"],
]
make_table(slide, table_data, Inches(0.8), Inches(1.8), Inches(7), row_height=Inches(0.5), font_size=14)

# Key differentiators
tf = add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.5))
tf_para(tf, "Key Technical Differentiators", size=20, bold=True, color=ACCENT_GREEN)
add_bullet(tf, "\u2022  Compound-grouped 5-fold cross-validation (no data leakage)", size=15, color=WHITE)
add_bullet(tf, "\u2022  Domain-validated: features align with known adhesive chemistry", size=15, color=WHITE)
add_bullet(tf, "     Epoxide ring breathing at 910\u2013915 cm\u207b\u00b9  |  Urethane N-H at 3300 cm\u207b\u00b9  |  Silicone Si-O-Si at 1000\u20131100 cm\u207b\u00b9", size=12, color=SUBTITLE_GRAY)
add_bullet(tf, "\u2022  Fisher ratios > 23 for both IR and Raman (excellent class separation)", size=15, color=WHITE)
add_bullet(tf, "\u2022  Dataset: 1,500+ curated adhesive spectra across 7 classes", size=15, color=WHITE)
add_bullet(tf, "\u2022  NMR excluded from production model (non-discriminative, Fisher ratio 0.02)", size=15, color=WHITE)


# ============================================================
# SLIDE 5: MARKET OPPORTUNITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
tf_para(tf, "MARKET OPPORTUNITY", size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.LEFT)
tf_para(tf, "$500\u2013620M TAM \u2014 No Direct Competitor", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

# TAM/SAM/SOM table
market_data = [
    ["Market Layer", "Value"],
    ["Total Addressable Market (TAM)", "$500\u2013620M"],
    ["Serviceable Addressable Market (SAM)", "~$165M"],
    ["Serviceable Obtainable Market (SOM, Yr 5)", "~$12M ARR"],
]
make_table(slide, market_data, Inches(0.8), Inches(1.8), Inches(6), row_height=Inches(0.45), font_size=13)

# Bottom-up
tf = add_text_box(slide, Inches(7.5), Inches(1.8), Inches(5.3), Inches(2.5))
tf_para(tf, "Bottom-Up Construction", size=18, bold=True, color=ACCENT_GREEN)
add_bullet(tf, "~2,100 companies with adhesive R&D programs", size=14, color=WHITE)
add_bullet(tf, "Tier 1 (>$10B): ~20 companies x $800K = $16M", size=14, color=WHITE)
add_bullet(tf, "Tier 2 ($1B\u201310B): ~180 co. x $400K = $72M", size=14, color=WHITE)
add_bullet(tf, "Tier 3 ($500M\u20131B): ~1,900 co. x $175K = $332M", size=14, color=WHITE)

# Market tailwinds
tf = add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(3))
tf_para(tf, "Market Tailwinds", size=18, bold=True, color=GOLD)
add_bullet(tf, "\u2022  PFAS reformulation wave \u2014 urgent need for rapid adhesive re-characterization", size=15, color=WHITE)
add_bullet(tf, "\u2022  EU Digital Product Passport regulations requiring material traceability", size=15, color=WHITE)
add_bullet(tf, "\u2022  AI-assisted formulation trend across specialty chemicals (5.7% CAGR)", size=15, color=WHITE)

# Key insight box
box = add_shape_bg(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), RGBColor(0x1A, 0x25, 0x3A))
tf2 = box.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "We are the CATEGORY CREATOR \u2014 $1.2M+ barrier to replicate our dataset and models"
p.font.size = Pt(18)
p.font.color.rgb = GOLD
p.font.bold = True
p.alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 6: BUSINESS MODEL & PRICING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_GREEN)

tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
tf_para(tf, "BUSINESS MODEL", size=14, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.LEFT)
tf_para(tf, "SaaS-First, Multi-Tier Revenue Model", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

# Pricing table
pricing_data = [
    ["Tier", "Price", "Target", "% Revenue"],
    ["SaaS Platform", "$2,500\u20138,000/mo", "R&D teams", "60%"],
    ["Per-Test API", "$75\u2013150/test", "QC labs", "25%"],
    ["Enterprise License", "$120\u2013250K/yr", "On-prem, dedicated", "15%"],
    ["Pilot Program", "$50K flat (6 mo)", "New customer POV", "Entry"],
    ["Academic", "$5\u201315K/yr", "Universities", "Self-serve"],
]
make_table(slide, pricing_data, Inches(0.8), Inches(1.7), Inches(11.5), row_height=Inches(0.4), font_size=12)

# Revenue build
revenue_data = [
    ["Year", "2026", "2027", "2028", "2029", "2030"],
    ["Customers", "24", "60", "100", "143", "187"],
    ["ARR", "$727K", "$2.2M", "$4.0M", "$6.4M", "$8.7M"],
    ["YoY Growth", "\u2014", "199%", "85%", "59%", "37%"],
]
make_table(slide, revenue_data, Inches(0.8), Inches(4.3), Inches(11.5), row_height=Inches(0.4), font_size=12)

tf = add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8))
tf_para(tf, "Per-test upside: 100K tests/year across 50 QC labs = $10\u201325M opportunity", size=16, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: FINANCIAL PROJECTIONS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
tf_para(tf, "FINANCIAL PROJECTIONS", size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.LEFT)
tf_para(tf, "Path to $8.7M ARR and 45% Net Margins by Year 5", size=32, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

fin_data = [
    ["Metric", "Year 1", "Year 2", "Year 3", "Year 4", "Year 5"],
    ["Revenue", "$727K", "$2.2M", "$4.0M", "$6.4M", "$8.7M"],
    ["Gross Margin", "58%", "68%", "74%", "77%", "78%"],
    ["EBITDA", "($523K)", "$274K", "$1.5M", "$3.2M", "$5.0M"],
    ["Net Margin", "(77%)", "8%", "27%", "39%", "45%"],
    ["Customers", "24", "60", "100", "143", "187"],
]
make_table(slide, fin_data, Inches(0.8), Inches(1.7), Inches(11.5), row_height=Inches(0.42), font_size=13)

# Unit economics
unit_data = [
    ["Metric", "Year 1", "Year 5"],
    ["CAC (fully loaded)", "$18.8K", "$8.5K"],
    ["LTV (3-year, blended)", "$62K", "$102K"],
    ["LTV:CAC", "3.3x", "12.0x"],
    ["Payback period", "14 months", "6 months"],
    ["Annual churn", "15%", "7%"],
]
make_table(slide, unit_data, Inches(0.8), Inches(4.3), Inches(6), row_height=Inches(0.38), font_size=12)

# Break-even callout
box = add_shape_bg(slide, Inches(7.5), Inches(4.3), Inches(5), Inches(2), RGBColor(0x1A, 0x25, 0x3A))
tf2 = box.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "BREAK-EVEN"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

pp = tf2.add_paragraph()
pp.text = "Month 26"
pp.font.size = Pt(36)
pp.font.color.rgb = GOLD
pp.font.bold = True
pp.alignment = PP_ALIGN.CENTER

pp2 = tf2.add_paragraph()
pp2.text = "~55 customers  |  $210K MRR"
pp2.font.size = Pt(14)
pp2.font.color.rgb = WHITE
pp2.alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 8: COMPETITIVE ADVANTAGES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
tf_para(tf, "COMPETITIVE ADVANTAGES", size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.LEFT)
tf_para(tf, "Why K-Dense Wins", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

comp_data = [
    ["Advantage", "K-Dense", "Traditional Labs", "DIY (Python)"],
    ["Speed", "Seconds", "Hours\u2013Days", "Minutes"],
    ["Accuracy", "95\u2013100%", "Variable", "Untested"],
    ["Sample Prep", "None", "Extensive", "N/A"],
    ["Cost per Test", "Low (SaaS)", "$50\u2013200+", "Engineer time"],
    ["Non-Destructive", "Yes", "No", "N/A"],
    ["7-Class Coverage", "Yes", "Yes", "Limited"],
]
make_table(slide, comp_data, Inches(0.8), Inches(1.7), Inches(11.5), row_height=Inches(0.4), font_size=12)

# Moats
tf = add_text_box(slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(2.5))
tf_para(tf, "Defensible Moats (Layered)", size=20, bold=True, color=ACCENT_GREEN)
moats = [
    "Proprietary dataset \u2014 1,500+ labeled adhesive spectra ($1.2M+ to replicate)",
    "Domain-validated models \u2014 chemically interpretable features confirmed by experts",
    "Compound-grouped validation \u2014 methodology exceeding industry standards",
    "Customer data moat \u2014 switching costs grow as customers index proprietary libraries",
    "First-mover advantage \u2014 no direct competitor in AI-powered adhesive classification",
]
for i, moat in enumerate(moats):
    add_bullet(tf, f"{i+1}.  {moat}", size=14, color=WHITE)


# ============================================================
# SLIDE 9: TEAM & CAPABILITIES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_GREEN)

tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
tf_para(tf, "TEAM", size=14, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.LEFT)
tf_para(tf, "World-Class Interdisciplinary Team", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

team = [
    ("Dr. Elena Vasquez", "CEO", "Strategic vision and company leadership"),
    ("Dr. Marcus Chen", "Chief Science Officer", "Research coordination across all departments"),
    ("Dr. Alexander Petrov", "Machine Learning Lead", "Model architecture, training, production deployment"),
    ("Dr. Nikolai Volkov", "Physical Sciences Lead", "Spectroscopy expertise & domain validation"),
    ("Dr. Priya Sharma", "Drug Discovery Lead", "Chemical validation & domain review"),
    ("Dr. Lin Wei", "Data Engineering Lead", "Data curation, pipeline, quality assurance"),
    ("Dr. Rosa Martinez", "Data Visualization", "Publication-quality figures & communication"),
    ("Dr. Paolo Ricci", "Market Research", "Market analysis & competitive intelligence"),
    ("Dr. Sarah Nakamura", "Financial Modeling", "Revenue projections & business planning"),
]

team_data = [["Name", "Role", "Focus"]] + [[n, r, f] for n, r, f in team]
make_table(slide, team_data, Inches(0.8), Inches(1.7), Inches(11.5), row_height=Inches(0.42), font_size=11)

tf = add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8))
tf_para(tf, "15+ specialists across ML, chemistry, data science, and business strategy", size=18, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10: GO-TO-MARKET STRATEGY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
tf_para(tf, "GO-TO-MARKET", size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.LEFT)
tf_para(tf, "Land-and-Expand with Design Partners", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

# Design partners
partners_data = [
    ["Target", "Pilot Prob.", "Rationale"],
    ["Evonik Industries", "88%", "Open Innovation (Creavis), digital chemistry"],
    ["Henkel AG", "82%", "World's largest adhesive mfr, Digital R&D"],
    ["Covestro AG", "75%", "Strong digital culture, PU adhesive focus"],
    ["Arkema/Bostik", "70%", "Post-acquisition data integration need"],
    ["3M Company", "64%", "Largest adhesive portfolio"],
]
make_table(slide, partners_data, Inches(0.8), Inches(1.7), Inches(6.5), row_height=Inches(0.4), font_size=12)

# Phases
phases_data = [
    ["Phase", "Timeline", "Revenue Target"],
    ["Foundation", "Q1\u2013Q2 2026", "\u2014"],
    ["Design Partners", "Q3\u2013Q4 2026", "$100\u2013150K"],
    ["Enterprise Expansion", "2027", "$1.5\u20132.5M ARR"],
    ["Scale", "2028\u20132030", "$5\u201312M ARR"],
]
make_table(slide, phases_data, Inches(8), Inches(1.7), Inches(4.8), row_height=Inches(0.4), font_size=12)

# GTM advantages
tf = add_text_box(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.5))
tf_para(tf, "Key GTM Advantages", size=18, bold=True, color=ACCENT_GREEN)
add_bullet(tf, "\u2022  95\u2013100% accuracy verifiable in a 15-minute live demo with customer spectra", size=15, color=WHITE)
add_bullet(tf, "\u2022  PFAS Reformulation Challenge \u2014 free public dataset drives 5\u201315 inbound enterprise inquiries", size=15, color=WHITE)
add_bullet(tf, "\u2022  DREAMS research paper in ACS Applied Materials & Interfaces builds scientific credibility", size=15, color=WHITE)


# ============================================================
# SLIDE 11: INVESTMENT ASK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), GOLD)

tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
tf_para(tf, "INVESTMENT ASK", size=14, bold=True, color=GOLD, alignment=PP_ALIGN.LEFT)
tf_para(tf, "$1.5M Seed  \u2192  $5M Series A", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

# Seed
seed_data = [
    ["Use of Funds (Seed $1.5M)", "Amount"],
    ["Product development & model optimization", "$600K"],
    ["Initial sales team (2 AEs + SDR)", "$350K"],
    ["Cloud infrastructure & hosting", "$200K"],
    ["Regulatory & compliance (SOC 2, ISO)", "$100K"],
    ["Working capital", "$250K"],
]
make_table(slide, seed_data, Inches(0.8), Inches(1.7), Inches(5.5), row_height=Inches(0.38), font_size=12)

# Series A
sa_data = [
    ["Use of Funds (Series A $5M)", "Amount"],
    ["Scale engineering (8\u219212 FTEs)", "$2.0M"],
    ["Expand sales & marketing", "$1.5M"],
    ["Enterprise on-prem features", "$800K"],
    ["International market entry", "$500K"],
    ["Working capital", "$200K"],
]
make_table(slide, sa_data, Inches(7), Inches(1.7), Inches(5.5), row_height=Inches(0.38), font_size=12)

# Milestones
tf = add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.5))
tf_para(tf, "Seed Milestones (18 months)", size=16, bold=True, color=ACCENT_GREEN)
add_bullet(tf, "\u2022  Launch commercial SaaS (Month 6)", size=13, color=WHITE)
add_bullet(tf, "\u2022  15+ paying customers (Month 12)", size=13, color=WHITE)
add_bullet(tf, "\u2022  $50K MRR (Month 15)", size=13, color=WHITE)
add_bullet(tf, "\u2022  Validate enterprise licensing (Month 18)", size=13, color=WHITE)

tf = add_text_box(slide, Inches(7), Inches(4.5), Inches(5.5), Inches(2.5))
tf_para(tf, "Series A Milestones", size=16, bold=True, color=ACCENT_GREEN)
add_bullet(tf, "\u2022  $200K MRR (Month 24)", size=13, color=WHITE)
add_bullet(tf, "\u2022  60+ customers (Month 30)", size=13, color=WHITE)
add_bullet(tf, "\u2022  Positive EBITDA (Month 26)", size=13, color=WHITE)
add_bullet(tf, "\u2022  First enterprise on-prem deployment", size=13, color=WHITE)


# ============================================================
# SLIDE 12: VISION / CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_GREEN)

tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
tf_para(tf, "VISION", size=14, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.LEFT)
tf_para(tf, "From Lab to Industry Standard", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

# Quote
box = add_shape_bg(slide, Inches(1.5), Inches(1.7), Inches(10), Inches(1), RGBColor(0x1A, 0x25, 0x3A))
tf2 = box.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = '"Every adhesive bond in every product deserves to be verified \u2014 instantly, accurately, and affordably."'
p.font.size = Pt(18)
p.font.color.rgb = GOLD
p.font.italic = True
p.alignment = PP_ALIGN.CENTER

# Roadmap
tf = add_text_box(slide, Inches(0.8), Inches(3), Inches(11.5), Inches(2))
add_bullet(tf, "Near-term: Production-ready adhesive classifier with 95\u2013100% accuracy", size=17, color=WHITE, bold=True)
add_bullet(tf, "Medium-term: Industry-standard AI testing platform integrated with major spectrometer brands", size=17, color=WHITE, bold=True)
add_bullet(tf, "Long-term: Universal spectral intelligence platform for material identification across industries", size=17, color=WHITE, bold=True)

# Scenario table
scenario_data = [
    ["Scenario", "Year 5 ARR", "Assumption"],
    ["Bear", "$1.5M", "Pilot struggles; build-in-house risk"],
    ["Base", "$8.7M", "3 pilots in 2026; steady expansion"],
    ["Bull", "$14M", "PFAS mandate urgency + adjacent verticals"],
]
make_table(slide, scenario_data, Inches(2), Inches(5), Inches(9), row_height=Inches(0.42), font_size=13)

# Contact
tf = add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.8))
tf_para(tf, "K-Dense Science Lab  \u2014  DREAMS Project  \u2014  Contact: Dr. Elena Vasquez, CEO", size=18, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "K-Dense_Investor_Pitch_Deck_March_2026.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
